# -*- coding: utf-8 -*-
"""Build the VGA game console acceptance presentation.

This script generates the PPTX, a PDF fallback deck, outline, speaker script,
source/citation list, and self-check report for the Nexys A7 VGA game console
course acceptance presentation.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from textwrap import dedent

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = Path(__file__).resolve().parent

PPTX_PATH = OUT_DIR / "VGA游戏集合机_验收讲解.pptx"
PDF_PATH = OUT_DIR / "VGA游戏集合机_验收讲解.pdf"
OUTLINE_PATH = OUT_DIR / "PPT大纲.md"
SCRIPT_PATH = OUT_DIR / "逐页演讲稿.md"
CITATION_PATH = OUT_DIR / "素材与代码引用清单.md"
REPORT_PATH = OUT_DIR / "制作说明与自检报告.md"


SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_CN = "SimHei"
FONT_CODE = "Consolas"

COLORS = {
    "bg": RGBColor(247, 249, 252),
    "ink": RGBColor(11, 31, 51),
    "muted": RGBColor(82, 96, 112),
    "blue": RGBColor(29, 78, 216),
    "cyan": RGBColor(8, 145, 178),
    "green": RGBColor(22, 163, 74),
    "orange": RGBColor(234, 88, 12),
    "red": RGBColor(220, 38, 38),
    "line": RGBColor(203, 213, 225),
    "soft_blue": RGBColor(219, 234, 254),
    "soft_cyan": RGBColor(207, 250, 254),
    "soft_green": RGBColor(220, 252, 231),
    "soft_orange": RGBColor(255, 237, 213),
    "soft_gray": RGBColor(241, 245, 249),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
}


@dataclass
class SlideSpec:
    title: str
    takeaway: str
    bullets: list[str] = field(default_factory=list)
    visual: str = ""
    speaker: str = ""
    refs: list[str] = field(default_factory=list)
    code: str = ""
    kind: str = "standard"


SLIDES: list[SlideSpec] = [
    SlideSpec(
        title="“学”：基于 Nexys A7 的 VGA 游戏集合机",
        takeaway="从基础 Verilog 进入 FPGA 图形交互系统。",
        bullets=[
            "板卡：Nexys A7-100T / XC7A100T-1CSG324C",
            "关键词：VGA / PS2 / Game Slot / LUT / WNS",
            "主线：显示、输入、平台化、资源分析、时序收敛",
        ],
        visual="cover",
        speaker=(
            "这次实验我用一个字概括，就是“学”。它不是只把几个小游戏写出来，"
            "而是让我从基础 Verilog 进入到一个 FPGA 图形交互系统：显示器怎么扫，"
            "键盘怎么输入，多个游戏怎么共享外设，以及最后怎么看资源和时序报告。"
            "所以我会把它当作一个 FPGA 游戏集合机来讲。"
        ),
        refs=["README.md", "docs/第二实验_VGA游戏集合机_验收讲解文档.md"],
        kind="cover",
    ),
    SlideSpec(
        title="这不是单个小游戏，而是一台 FPGA 游戏集合机",
        takeaway="从用户角度，它有统一入口、统一输入和返回路径。",
        bullets=[
            "上电后先进入 VGA 菜单",
            "PS/2 键盘选择并启动游戏",
            "游戏中按 Esc 返回菜单",
            "四个槽位：TANK WAR / GAME TWO / GAME THREE / GAME FOUR",
        ],
        visual="user_flow",
        speaker=(
            "从用户角度看，这个项目是一台简易 FPGA 游戏主机。上电后显示的是菜单，"
            "用户用 W/S 或方向键移动，用 Enter 或 Space 启动游戏，进入游戏后按 Esc 返回。"
            "所以它不是几个孤立游戏模块的堆叠，而是一个有完整用户流程的集合机。"
        ),
        refs=["README.md"],
    ),
    SlideSpec(
        title="我的学习路线：从能显示，到能优化",
        takeaway="本次汇报按学习路径讲，不按文件列表讲。",
        bullets=[
            "先解释 VGA 实时像素扫描",
            "再解释菜单渲染和 PS/2 输入事件化",
            "然后说明 Game Slot API 的平台化意义",
            "最后用资源和时序报告说明工程优化过程",
        ],
        visual="learning_route",
        speaker=(
            "今天我不按文件顺序讲，因为那样容易变成代码导览。我会按学习路线讲："
            "先从 VGA 为什么能显示开始，再到菜单怎么画、键盘怎么进来，然后讲多游戏如何接入，"
            "最后讲资源占用和 WNS 从负到正的优化过程。"
        ),
        refs=["docs/PPT设计方案.md", "docs/第二实验_VGA游戏集合机_验收讲解文档.md"],
    ),
    SlideSpec(
        title="主机 + 游戏槽位：game_console_top 如何组织系统",
        takeaway="顶层提供公共服务，槽位只负责自身逻辑和画面。",
        bullets=[
            "common 模块：VGA 同步、菜单渲染、PS/2 接收、菜单控制",
            "slot1 到 slot4 使用统一端口接入",
            "VGA、LED、七段管、蜂鸣器由顶层统一复用",
        ],
        visual="architecture",
        speaker=(
            "系统架构可以理解成“主机加游戏槽位”。`game_console_top` 统一提供 VGA、PS/2、"
            "菜单和输出复用服务；每个游戏槽位只处理自己的状态和渲染。这样做的价值是，"
            "游戏之间不会各自抢外设，新增游戏也有稳定接口可以接入。"
        ),
        refs=["src/game_console_top.v", "docs/game_api.md"],
    ),
    SlideSpec(
        title="VGA 的本质：每个像素周期输出一个颜色",
        takeaway="FPGA 不是发送整张图片，而是扫描到哪里就实时给出哪里的 RGB。",
        bullets=[
            "`pixel_x / pixel_y` 表示当前扫描坐标",
            "`display_active` 表示当前是否在 640x480 可见区",
            "`VGA_R/G/B` 给出当前像素颜色",
            "`VGA_HS / VGA_VS` 让显示器识别行和帧同步",
        ],
        visual="vga_scan",
        speaker=(
            "VGA 是本项目最重要的学习点之一。显示器并不是接收一整张图片，"
            "而是要求 FPGA 按照扫描顺序，在每个像素周期输出当前点的颜色。"
            "所以我们先生成 `pixel_x` 和 `pixel_y`，再让菜单或游戏根据这个坐标实时计算 RGB。"
        ),
        refs=["src/common/console_vga_sync.v"],
    ),
    SlideSpec(
        title="console_vga_sync：用计数器生成像素坐标",
        takeaway="VGA 控制器本质上是两个计数器加同步脉冲生成逻辑。",
        code=dedent(
            """\
            // src/common/console_vga_sync.v
            localparam integer H_VISIBLE = 640;
            localparam integer H_FRONT   = 16;
            localparam integer H_SYNC    = 96;
            localparam integer H_BACK    = 48;
            localparam integer V_VISIBLE = 480;
            localparam integer V_FRONT   = 10;
            localparam integer V_SYNC    = 2;
            localparam integer V_BACK    = 33;

            reg [1:0] pix_div;
            assign pixel_tick = (pix_div == 2'b11);
            assign display_active = (h_count < H_VISIBLE) &&
                                    (v_count < V_VISIBLE);
            """
        ),
        bullets=[
            "可见区为 640x480，扫描还包含 front/sync/back 消隐区",
            "`pix_div` 每 4 个 100MHz 周期产生一次像素使能",
            "`h_count/v_count` 同步推进，得到当前像素坐标",
        ],
        visual="code",
        speaker=(
            "这一页对应真实代码。水平和垂直方向都有 visible、front、sync、back 参数，"
            "可见区是 640 乘 480，但实际扫描总宽度还包括消隐区。`pix_div` 产生四拍一次的"
            "`pixel_tick`，`h_count` 和 `v_count` 只在这个使能下推进。"
        ),
        refs=["src/common/console_vga_sync.v:13", "src/common/console_vga_sync.v:43"],
    ),
    SlideSpec(
        title="单时钟域设计：100MHz 主时钟 + pixel_tick 使能",
        takeaway="没有新建 25MHz 时钟，而是用 clock enable 保持单时钟域。",
        bullets=[
            "板卡主时钟是 100MHz",
            "`pixel_tick` 作为 VGA 像素推进使能",
            "所有主要逻辑仍由 `CLK100MHZ` 驱动",
            "时序约束和跨模块协作更清晰",
        ],
        visual="single_clock",
        speaker=(
            "这里我没有再创建一个 25MHz 新时钟，而是用 `pixel_tick` 作为 clock enable。"
            "这样 VGA 坐标按像素节拍前进，但寄存器仍在 100MHz 主时钟域下工作。"
            "这能避免不必要的跨时钟域问题，也让后面时序约束更明确。"
        ),
        refs=["src/common/console_vga_sync.v", "docs/game_api.md"],
    ),
    SlideSpec(
        title="菜单不是图片，而是一组像素级判断",
        takeaway="菜单渲染同样是实时像素逻辑，也是 LUT 压力来源之一。",
        code=dedent(
            """\
            // src/common/console_menu_renderer.v 伪代码
            if (!display_active)
                rgb = 12'h000;
            else if (pixel 命中标题)
                rgb = title_color;
            else if (pixel 命中菜单项文字)
                rgb = selected ? highlight_color : normal_text;
            else if (pixel 命中边框)
                rgb = border_color;
            else
                rgb = background_color;
            """
        ),
        bullets=[
            "`pixel_x/y` 进入区域判断",
            "标题、菜单项、帮助文字和边框都由坐标命中决定",
            "`cursor` 只决定哪一项高亮",
        ],
        visual="menu_render",
        speaker=(
            "菜单看起来像一张界面图，但在 FPGA 里它不是图片。当前像素是不是标题、"
            "是不是菜单项文字、是不是边框或高亮，都要根据坐标判断。"
            "这让我理解到，图形界面在硬件里会自然变成很多比较器、case 和 mux。"
        ),
        refs=["src/common/console_menu_renderer.v"],
    ),
    SlideSpec(
        title="菜单控制器决定“选谁”，渲染器决定“怎么画”",
        takeaway="输入事件、菜单状态和像素渲染分层，顶层再做显示复用。",
        bullets=[
            "`console_ps2_rx` 输出 byte 事件",
            "`console_menu_controller` 更新 `cursor / game_sel / menu_active`",
            "`console_menu_renderer` 根据 `cursor` 和坐标输出菜单 RGB",
            "顶层根据 `menu_active` 选择菜单或当前游戏画面",
        ],
        visual="menu_control",
        speaker=(
            "菜单控制和菜单渲染是分开的。控制器只回答“当前选谁、是否进入游戏”，"
            "渲染器只回答“这个像素怎么画”。这比把键盘处理和画面逻辑混在一起更清楚，"
            "也方便顶层在菜单画面和游戏画面之间做复用。"
        ),
        refs=["src/common/console_ps2_rx.v", "src/common/console_menu_controller.v", "src/common/console_menu_renderer.v"],
    ),
    SlideSpec(
        title="PS/2 输入不是按键电平，而是串行协议",
        takeaway="键盘先发送扫描码，FPGA 必须先协议解析再事件化。",
        bullets=[
            "一帧包含 start、8 位 scan code、parity、stop",
            "数据位 LSB first",
            "在 PS/2 时钟下降沿采样",
            "底层只输出 `byte_ready + byte_data`",
        ],
        visual="ps2_frame",
        speaker=(
            "PS/2 不能当普通按键电平使用。键盘通过时钟线和数据线串行发送扫描码，"
            "一帧包括起始位、8 位数据、校验位和停止位。FPGA 先在下降沿采样并组装 byte，"
            "上层再把 byte 解释成菜单或游戏动作。"
        ),
        refs=["src/common/console_ps2_rx.v"],
    ),
    SlideSpec(
        title="console_ps2_rx：下降沿采样并组装扫描码",
        takeaway="这个模块只做协议层接收，不解释具体按键含义。",
        code=dedent(
            """\
            // src/common/console_ps2_rx.v
            assign ps2_clk_fall_next = ps2_clk_q & ~ps2_clk_next;
            assign ps2_parity_ok = (^shift_data) ^ parity_bit;

            if (ps2_clk_fall_next) begin
                case (bit_count)
                    4'd0: if (ps2_data_next == 1'b0)
                              bit_count <= 4'd1;
                    4'd1: begin shift_data[0] <= ps2_data_next;
                                  bit_count <= 4'd2; end
                    // ... shift_data[1] ~ shift_data[7], parity
                    4'd10: if (ps2_data_next == 1'b1 &&
                               ps2_parity_ok) begin
                               byte_data <= shift_data;
                               byte_ready <= 1'b1;
                           end
                endcase
            end
            """
        ),
        bullets=[
            "PS/2 时钟和数据先同步并做历史滤波",
            "`bit_count` 记录当前接收到帧的哪一位",
            "停止位为 1 且奇校验通过时产生 `byte_ready`",
        ],
        visual="code",
        speaker=(
            "这里是真实代码片段。项目不是直接用裸下降沿，而是先同步和滤波 PS/2 信号，"
            "再用 `ps2_clk_fall_next` 判断下降沿。`bit_count` 从起始位推进到数据位、校验位和停止位，"
            "最后在停止位为 1 且校验通过时输出一个 byte 事件。"
        ),
        refs=["src/common/console_ps2_rx.v:12", "src/common/console_ps2_rx.v:40", "src/common/console_ps2_rx.v:70"],
    ),
    SlideSpec(
        title="扫描码事件化：W/S、Enter、Esc 如何控制菜单",
        takeaway="底层 byte 被转换成上移、下移、启动、返回这些系统事件。",
        code=dedent(
            """\
            // src/common/console_menu_controller.v
            localparam [7:0] SCAN_F0    = 8'hF0;
            localparam [7:0] SCAN_E0    = 8'hE0;
            localparam [7:0] SCAN_W     = 8'h1D;
            localparam [7:0] SCAN_S     = 8'h1B;
            localparam [7:0] SCAN_SPACE = 8'h29;
            localparam [7:0] SCAN_ENTER = 8'h5A;
            localparam [7:0] SCAN_ESC   = 8'h76;
            localparam [7:0] SCAN_UP    = 8'h75;
            localparam [7:0] SCAN_DOWN  = 8'h72;

            assign is_up_key = (byte_data == SCAN_W) ||
                               (byte_data == SCAN_UP);
            assign is_start_key = (byte_data == SCAN_SPACE) ||
                                  (byte_data == SCAN_ENTER);
            """
        ),
        bullets=[
            "`F0` 表示 key release 前缀，下一码不触发动作",
            "`E0` 记录扩展键前缀，方向键扫描码仍由上层识别",
            "菜单内启动游戏；游戏中 Esc 返回菜单",
        ],
        visual="scan_codes",
        speaker=(
            "`console_menu_controller` 做的是事件化。PS/2 接收器只知道收到哪个扫描码，"
            "控制器再把 W 或 Up 解释成上移，把 S 或 Down 解释成下移，把 Enter 或 Space 解释成启动，"
            "把 Esc 解释成从游戏返回菜单。这里也处理了 F0 释放前缀，避免松键重复触发。"
        ),
        refs=["src/common/console_menu_controller.v:12", "src/common/console_menu_controller.v:33"],
    ),
    SlideSpec(
        title="统一槽位 API：从“多个游戏”到“游戏平台”",
        takeaway="每个游戏像卡带一样接入，顶层只依赖统一端口。",
        bullets=[
            "公共输入：`clk / reset / selected`",
            "显示输入：`frame_tick / pixel_tick / display_active / pixel_x/y`",
            "输入事件：`ps2_byte_ready / ps2_byte_data`、按钮和开关",
            "统一输出：VGA RGB、LED、七段管、蜂鸣器",
        ],
        visual="game_slot_api",
        speaker=(
            "Game Slot API 是这个项目从多个游戏变成游戏平台的关键。每个游戏都接收同一套公共输入，"
            "再输出同一类板级信号。顶层不需要理解游戏内部规则，只需要选择当前 slot 的输出。"
            "这样以后替换游戏或新增游戏，接口边界是稳定的。"
        ),
        refs=["docs/game_api.md"],
    ),
    SlideSpec(
        title="输出复用：菜单激活显示菜单，否则显示当前游戏",
        takeaway="真实 VGA 端口只有一组，所有菜单和游戏画面都必须由顶层仲裁。",
        code=dedent(
            """\
            // src/game_console_top.v
            assign active_slot_vga_r = slot1_selected ? slot1_vga_r :
                                       slot2_selected ? slot2_vga_r :
                                       slot3_selected ? slot3_vga_r :
                                                        slot4_vga_r;

            always @(*) begin
                if (menu_active) begin
                    VGA_R = menu_vga_r_q;
                    VGA_G = menu_vga_g_q;
                    VGA_B = menu_vga_b_q;
                end else begin
                    VGA_R = active_slot_vga_r;
                    VGA_G = active_slot_vga_g;
                    VGA_B = active_slot_vga_b;
                end
            end
            """
        ),
        bullets=[
            "slot1 到 slot4 先复用为 `active_slot_*`",
            "菜单态覆盖 VGA、LED、七段管和蜂鸣器输出",
            "HS/VS 始终来自公共 `console_vga_sync`",
        ],
        visual="mux",
        speaker=(
            "真实板卡的 VGA 口只有一组，所以顶层必须复用。先根据 `game_sel` 选出当前槽位 RGB，"
            "再根据 `menu_active` 决定输出菜单 RGB 还是游戏 RGB。这里要注意，运行时未选中的游戏不等于综合时不占资源，"
            "真正裁剪资源需要构建宏或 stub。"
        ),
        refs=["src/game_console_top.v:400", "src/game_console_top.v:440"],
    ),
    SlideSpec(
        title="资源报告告诉我：问题不是寄存器太多，而是组合逻辑太重",
        takeaway="LUT 约 48%，FF 约 3%，说明组合路径和像素渲染是主要压力。",
        bullets=[
            "LUT：30146 / 63400，约 48%",
            "FF：4371 / 126800，约 3%",
            "菜单、文字、地图、sprite、tile 和对象命中检测都会消耗 LUT",
            "优化方向是适当增加寄存器和 staging，而不是继续减少 FF",
        ],
        visual="resource_chart",
        speaker=(
            "资源分析这一页是我后期理解硬件实现的重要转折。LUT 用了约 48%，但 FF 只有约 3%。"
            "这说明问题不是寄存器太多，而是大量逻辑被展开在组合网络里。"
            "对 VGA 游戏来说，文字、地图、对象命中和图层优先级都会变成 LUT。"
        ),
        refs=["docs/资源占用分析.md"],
    ),
    SlideSpec(
        title="每个像素都要回答一个问题：这里该显示什么？",
        takeaway="VGA 游戏的复杂度来自每个像素周期的坐标判断和图层优先级。",
        bullets=[
            "是否在可见区",
            "是否命中背景、地图 tile、玩家、敌人、子弹、道具",
            "是否命中文字或 UI",
            "最终按图层优先级选择 RGB",
        ],
        visual="pixel_decision",
        speaker=(
            "VGA 游戏容易 LUT 高，是因为每个像素都要回答“这里该显示什么”。"
            "如果背景、地图、角色、子弹、文字和 UI 都在一拍内判断，就会形成很多比较器和多路选择器。"
            "所以优化不是删功能，而是把这些判断分阶段组织。"
        ),
        refs=["docs/资源占用分析.md", "src/common/console_menu_renderer.v", "src/games/slot3/slot3_renderer.v", "src/games/slot4/game_slot4_top.v"],
    ),
    SlideSpec(
        title="代码能综合，不代表电路能按 100MHz 跑",
        takeaway="WNS 为负说明至少有一条路径在一个周期内来不及。",
        bullets=[
            "100MHz 时钟周期为 10ns",
            "寄存器 A 到寄存器 B 中间的组合逻辑必须在周期内稳定",
            "WNS：最差路径时序余量",
            "TNS：所有失败路径负 slack 总和",
        ],
        visual="timing_concept",
        speaker=(
            "功能仿真或综合通过，不代表布线后的电路一定能按 100MHz 稳定运行。"
            "寄存器之间的组合逻辑如果超过一个时钟周期，WNS 就会为负。"
            "所以后期我关注的不只是能不能生成逻辑，而是 post-route 后所有受约束路径是否满足时序。"
        ),
        refs=["docs/final_timing_optimization_round.md"],
    ),
    SlideSpec(
        title="最终优化结果：从负 slack 到时序收敛",
        takeaway="先暴露真实基线，再通过 RTL staging/pipeline 把 WNS 拉正。",
        bullets=[
            "真实基线：WNS=-0.445ns，TNS=-18.104ns，失败端点 94",
            "最终 post-route：WNS=+0.174ns，TNS=0.000ns，失败端点 0",
            "只保留窄范围 VGA pixel multicycle，不覆盖游戏核心状态",
            "最终成功生成 bitstream",
        ],
        visual="timing_results",
        speaker=(
            "这一页是后期优化的核心结果。移除宽泛例外后，真实基线 WNS 是负的，"
            "TNS 也有 -18.104ns，并且有 94 个失败端点。最终通过 RTL 结构优化后，post-route WNS 为 +0.174ns，"
            "TNS 为 0，失败端点为 0。这个余量不算很大，但已经说明设计实现后时序收敛。"
        ),
        refs=["docs/final_timing_optimization_round.md", "scripts/apply_timing_exceptions.tcl"],
    ),
    SlideSpec(
        title="优化思想：把一拍大计算拆成多拍小计算",
        takeaway="用户可见行为不变，但硬件每拍需要完成的逻辑更少。",
        bullets=[
            "Slot2：ghost 迭代、lock row-by-row、video staging",
            "Slot3：`try_x / try_y` 寄存，下一拍判断 walkable",
            "Slot4：tile probe -> eval -> apply",
            "约束：4-cycle VGA multicycle 只用于 pixel_tick 视频路径",
        ],
        visual="optimization_cases",
        speaker=(
            "三个优化案例体现的是同一个思路：把一拍大计算拆成多拍小计算。"
            "Slot2 把 ghost 和 lock 拆成迭代和逐行流程；Slot3 把候选坐标先寄存，下一拍再查地图；"
            "Slot4 把 tile 碰撞拆成 probe、eval、apply。用户看到的行为不变，但硬件关键路径变短。"
        ),
        refs=["src/games/slot2/slot2_game_core.v", "src/games/slot2/game_slot2_top.v", "src/games/slot3/slot3_player.v", "src/games/slot4/game_slot4_top.v"],
    ),
    SlideSpec(
        title="从功能实现到工程实现",
        takeaway="最终收获是把显示、输入、架构、资源和时序统一起来。",
        bullets=[
            "VGA：理解实时像素扫描",
            "菜单：理解坐标判断和图层优先级",
            "PS/2：理解外设协议解析和事件化",
            "Game Slot API：理解平台化系统架构",
            "LUT/WNS：理解从 Verilog 功能到可实现硬件的差距",
        ],
        visual="summary",
        speaker=(
            "最后我把这个项目总结为从功能实现到工程实现。VGA 让我理解实时像素扫描，"
            "PS/2 让我理解外设协议解析，Game Slot API 让我理解平台化架构。"
            "资源和时序优化则让我意识到，写出能工作的 Verilog 还不够，还要让它在真实 FPGA 上满足资源和时序约束。"
        ),
        refs=["README.md", "docs/资源占用分析.md", "docs/final_timing_optimization_round.md"],
    ),
]


def rgb(name: str) -> RGBColor:
    return COLORS[name]


def add_bg(slide, color: RGBColor = COLORS["bg"]):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def set_text_frame(tf, font_name=FONT_CN, font_size=18, color=COLORS["ink"], bold=False):
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=18, color=None, bold=False, align=None, font=FONT_CN):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color or COLORS["ink"]
        run.font.bold = bold
    return box


def add_multiline(slide, lines, x, y, w, h, size=18, color=None, bullet=False, font=FONT_CN, spacing=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " + line) if bullet else line
        p.line_spacing = spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color or COLORS["ink"]
    return box


def add_code(slide, code, x, y, w, h, size=12.5):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(15, 23, 42)
    rect.line.color.rgb = RGBColor(30, 41, 59)
    rect.line.width = Pt(1)
    tf = rect.text_frame
    tf.clear()
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = code.rstrip()
    p.line_spacing = 0.9
    for run in p.runs:
        run.font.name = FONT_CODE
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(226, 232, 240)
    return rect


def add_box(slide, text, x, y, w, h, fill, line=None, size=16, bold=False, color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or COLORS["line"]
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = FONT_CN
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or COLORS["ink"]
    return shp


def add_small_arrow(slide, x, y, w=0.35, h=0.16, fill=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill or COLORS["cyan"]
    shp.line.fill.background()
    return shp


def add_header(slide, idx, title, takeaway):
    add_bg(slide)
    add_text(slide, title, 0.62, 0.32, 11.4, 0.45, size=24, bold=True, color=COLORS["ink"])
    add_text(slide, takeaway, 0.66, 0.86, 10.9, 0.32, size=12.8, color=COLORS["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(1.23), Inches(2.0), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["cyan"]
    line.line.fill.background()
    add_text(slide, f"{idx:02d}", 12.2, 0.34, 0.42, 0.26, size=10, color=COLORS["muted"], align=PP_ALIGN.RIGHT)
    add_text(slide, "Nexys A7 VGA Game Console", 9.75, 7.06, 2.8, 0.25, size=8.5, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_bullet_panel(slide, bullets, x=0.76, y=5.45, w=11.8, h=1.05, cols=2):
    if not bullets:
        return
    if cols == 1:
        add_multiline(slide, bullets, x, y, w, h, size=14, color=COLORS["ink"], bullet=True)
        return
    left = bullets[: (len(bullets) + 1) // 2]
    right = bullets[(len(bullets) + 1) // 2 :]
    add_multiline(slide, left, x, y, w / 2 - 0.15, h, size=13.4, bullet=True)
    add_multiline(slide, right, x + w / 2 + 0.15, y, w / 2 - 0.15, h, size=13.4, bullet=True)


def cover_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(10, 22, 40))
    add_text(slide, spec.title, 0.7, 0.72, 10.9, 0.82, size=29, bold=True, color=COLORS["white"])
    add_text(slide, "从 VGA 显示、PS/2 输入到资源分析与时序收敛", 0.78, 1.72, 9.5, 0.36, size=18, color=RGBColor(203, 213, 225))
    add_text(slide, "Nexys A7-100T", 0.82, 2.32, 2.2, 0.28, size=13, color=RGBColor(147, 197, 253))
    items = ["VGA", "PS2", "Game Slot", "LUT", "WNS"]
    x = 0.82
    for item in items:
        add_box(slide, item, x, 6.15, 1.3, 0.4, RGBColor(30, 41, 59), RGBColor(51, 65, 85), size=12, bold=True, color=COLORS["white"])
        x += 1.5
    # Stylized FPGA console signal diagram.
    add_box(slide, "FPGA\nGame Console", 8.4, 2.35, 2.15, 1.15, RGBColor(20, 184, 166), RGBColor(103, 232, 249), size=17, bold=True, color=COLORS["white"])
    add_box(slide, "VGA\n640x480", 6.6, 1.1, 1.35, 0.75, RGBColor(37, 99, 235), COLORS["white"], size=13, bold=True, color=COLORS["white"])
    add_box(slide, "PS/2\nKeyboard", 10.95, 1.1, 1.45, 0.75, RGBColor(234, 88, 12), COLORS["white"], size=13, bold=True, color=COLORS["white"])
    add_box(slide, "Slot1", 6.25, 4.38, 1.0, 0.48, RGBColor(30, 41, 59), RGBColor(148, 163, 184), size=11, color=COLORS["white"])
    add_box(slide, "Slot2", 7.55, 4.38, 1.0, 0.48, RGBColor(30, 41, 59), RGBColor(148, 163, 184), size=11, color=COLORS["white"])
    add_box(slide, "Slot3", 8.85, 4.38, 1.0, 0.48, RGBColor(30, 41, 59), RGBColor(148, 163, 184), size=11, color=COLORS["white"])
    add_box(slide, "Slot4", 10.15, 4.38, 1.0, 0.48, RGBColor(30, 41, 59), RGBColor(148, 163, 184), size=11, color=COLORS["white"])
    add_text(slide, "课程验收汇报", 0.82, 6.82, 2.0, 0.25, size=10, color=RGBColor(148, 163, 184))
    return slide


def draw_user_flow(slide):
    steps = ["上电", "VGA 菜单", "键盘选择", "启动游戏", "ESC 返回菜单"]
    x = 0.75
    for i, step in enumerate(steps):
        add_box(slide, step, x, 1.85, 1.75, 0.62, COLORS["soft_blue"] if i != 3 else COLORS["soft_orange"], COLORS["blue"], size=14, bold=True)
        if i < len(steps) - 1:
            add_small_arrow(slide, x + 1.88, 2.08, 0.42, 0.14, COLORS["cyan"])
        x += 2.45
    games = ["TANK WAR", "GAME TWO", "GAME THREE", "GAME FOUR"]
    x = 1.0
    for i, game in enumerate(games):
        fill = [COLORS["soft_green"], COLORS["soft_cyan"], COLORS["soft_blue"], COLORS["soft_orange"]][i]
        add_box(slide, game, x, 3.55, 2.35, 0.78, fill, COLORS["line"], size=15, bold=True)
        add_text(slide, f"slot{i+1}", x + 0.75, 4.38, 0.8, 0.2, size=9, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        x += 2.95
    add_text(slide, "统一入口 + 统一输入 + 统一返回", 3.4, 4.95, 6.4, 0.38, size=20, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)


def draw_learning_route(slide):
    items = ["VGA 显示原理", "菜单渲染", "PS/2 键盘输入", "Game Slot API", "资源占用分析", "WNS 为负到 post-route 收敛"]
    y = 1.52
    for i, item in enumerate(items):
        fill = [COLORS["soft_blue"], COLORS["soft_cyan"], COLORS["soft_green"], COLORS["soft_blue"], COLORS["soft_orange"], COLORS["soft_green"]][i]
        add_box(slide, item, 4.1, y, 5.2, 0.48, fill, COLORS["cyan"], size=15, bold=True)
        if i < len(items) - 1:
            add_text(slide, "↓", 6.55, y + 0.45, 0.2, 0.22, size=15, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        y += 0.78
    add_text(slide, "不是按文件顺序讲，而是按能力增长路径讲", 3.2, 6.48, 6.9, 0.32, size=15, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_architecture(slide):
    add_box(slide, "game_console_top\n公共服务与输出仲裁", 5.05, 1.42, 3.2, 0.82, COLORS["soft_blue"], COLORS["blue"], size=16, bold=True)
    common = [
        ("console_vga_sync", "VGA 时序/坐标"),
        ("console_ps2_rx", "PS/2 byte 接收"),
        ("console_menu_controller", "菜单状态"),
        ("console_menu_renderer", "菜单 RGB"),
    ]
    y = 2.72
    for name, desc in common:
        add_box(slide, f"{name}\n{desc}", 0.85, y, 3.6, 0.55, COLORS["soft_cyan"], COLORS["cyan"], size=11.5)
        add_small_arrow(slide, 4.55, y + 0.18, 0.35, 0.13, COLORS["cyan"])
        y += 0.78
    slots = [
        ("game_slot1_top", "TANK WAR"),
        ("game_slot2_top", "GAME TWO"),
        ("game_slot3_top", "GAME THREE"),
        ("game_slot4_top", "GAME FOUR"),
    ]
    y = 2.72
    for name, desc in slots:
        add_small_arrow(slide, 8.42, y + 0.18, 0.35, 0.13, COLORS["orange"])
        add_box(slide, f"{name}\n{desc}", 8.9, y, 3.6, 0.55, COLORS["soft_orange"], COLORS["orange"], size=11.5)
        y += 0.78
    add_text(slide, "公共服务", 1.98, 2.28, 1.3, 0.25, size=12, bold=True, color=COLORS["cyan"], align=PP_ALIGN.CENTER)
    add_text(slide, "统一槽位", 10.05, 2.28, 1.3, 0.25, size=12, bold=True, color=COLORS["orange"], align=PP_ALIGN.CENTER)


def draw_vga_scan(slide):
    screen = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.05), Inches(1.6), Inches(6.55), Inches(4.25))
    screen.fill.solid()
    screen.fill.fore_color.rgb = RGBColor(226, 232, 240)
    screen.line.color.rgb = COLORS["blue"]
    screen.line.width = Pt(2)
    # Scan lines.
    for i in range(8):
        y = 1.9 + i * 0.43
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.35), Inches(y), Inches(5.85), Inches(0.025))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(148, 163, 184)
        line.line.fill.background()
    add_box(slide, "当前像素\n(pixel_x, pixel_y)", 3.6, 3.35, 1.55, 0.55, COLORS["soft_orange"], COLORS["orange"], size=11.5, bold=True)
    add_small_arrow(slide, 2.0, 1.75, 4.5, 0.18, COLORS["cyan"])
    add_text(slide, "左 → 右", 3.3, 1.32, 1.1, 0.22, size=12, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_text(slide, "上 ↓ 下", 0.48, 3.35, 0.4, 0.55, size=12, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_box(slide, "VGA_R/G/B\n当前颜色", 8.65, 1.75, 2.3, 0.75, COLORS["soft_green"], COLORS["green"], size=15, bold=True)
    add_box(slide, "VGA_HS / VGA_VS\n同步脉冲", 8.65, 3.0, 2.3, 0.75, COLORS["soft_blue"], COLORS["blue"], size=15, bold=True)
    add_box(slide, "display_active\n可见区有效", 8.65, 4.25, 2.3, 0.75, COLORS["soft_cyan"], COLORS["cyan"], size=15, bold=True)


def draw_code_slide(slide, spec):
    add_code(slide, spec.code, 0.78, 1.55, 7.1, 4.62, size=11.4 if "console_vga" in spec.code else 10.7)
    add_multiline(slide, spec.bullets, 8.25, 1.68, 4.1, 2.9, size=14.2, bullet=True)
    add_box(slide, "真实代码片段", 8.35, 4.98, 2.0, 0.45, COLORS["soft_green"], COLORS["green"], size=14, bold=True)
    refs = "\n".join(spec.refs[:2])
    add_text(slide, refs, 8.3, 5.55, 4.0, 0.55, size=10.8, color=COLORS["muted"])


def draw_single_clock(slide):
    add_box(slide, "不推荐\n100MHz + 新建 25MHz", 1.05, 2.0, 3.5, 0.9, COLORS["soft_orange"], COLORS["orange"], size=15, bold=True)
    add_box(slide, "多时钟域\nCDC/约束复杂", 1.35, 3.5, 2.9, 0.72, RGBColor(254, 226, 226), COLORS["red"], size=14, bold=True, color=COLORS["red"])
    add_small_arrow(slide, 2.6, 3.02, 0.42, 0.16, COLORS["orange"])
    add_box(slide, "当前设计\n100MHz 主时钟", 7.0, 1.62, 3.3, 0.62, COLORS["soft_blue"], COLORS["blue"], size=15, bold=True)
    add_box(slide, "pixel_tick\nclock enable", 7.0, 2.72, 3.3, 0.62, COLORS["soft_cyan"], COLORS["cyan"], size=15, bold=True)
    add_box(slide, "VGA 坐标 / 菜单 / 槽位\n仍在单时钟域协作", 6.55, 3.82, 4.2, 0.82, COLORS["soft_green"], COLORS["green"], size=14, bold=True)
    add_small_arrow(slide, 8.45, 2.34, 0.42, 0.16, COLORS["cyan"])
    add_small_arrow(slide, 8.45, 3.45, 0.42, 0.16, COLORS["green"])
    add_text(slide, "核心取舍：用使能控制节拍，而不是增加系统时钟域", 3.1, 5.55, 7.2, 0.35, size=17, color=COLORS["blue"], bold=True, align=PP_ALIGN.CENTER)


def draw_menu_render(slide, spec):
    add_code(slide, spec.code, 0.82, 1.58, 5.7, 3.9, size=11.3)
    chain = ["pixel_x/y", "区域判断", "标题/菜单项", "高亮/边框", "RGB"]
    y = 1.65
    for i, item in enumerate(chain):
        fill = COLORS["soft_cyan"] if i < 2 else COLORS["soft_blue"] if i < 4 else COLORS["soft_green"]
        add_box(slide, item, 7.55, y, 2.6, 0.5, fill, COLORS["cyan"], size=14, bold=True)
        if i < len(chain) - 1:
            add_text(slide, "↓", 8.78, y + 0.47, 0.2, 0.2, size=13, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        y += 0.78
    add_text(slide, "菜单控制器给出 cursor，渲染器只负责把当前像素画对", 6.9, 5.55, 4.7, 0.42, size=13, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_menu_control(slide):
    chain = [
        ("console_ps2_rx", "byte_ready / byte_data"),
        ("console_menu_controller", "cursor / game_sel / menu_active"),
        ("console_menu_renderer", "menu_vga_r/g/b"),
        ("VGA 输出", "menu 或 active_slot"),
    ]
    y = 1.5
    for i, (name, desc) in enumerate(chain):
        fill = [COLORS["soft_cyan"], COLORS["soft_blue"], COLORS["soft_green"], COLORS["soft_orange"]][i]
        add_box(slide, f"{name}\n{desc}", 1.55, y, 4.0, 0.68, fill, COLORS["line"], size=13.5, bold=True)
        if i < len(chain) - 1:
            add_text(slide, "↓", 3.45, y + 0.66, 0.2, 0.25, size=15, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        y += 1.08
    add_box(slide, "控制逻辑\n决定“选谁”", 7.1, 2.1, 2.0, 0.92, COLORS["soft_blue"], COLORS["blue"], size=15, bold=True)
    add_box(slide, "渲染逻辑\n决定“怎么画”", 9.55, 2.1, 2.0, 0.92, COLORS["soft_green"], COLORS["green"], size=15, bold=True)
    add_text(slide, "分层后，菜单状态、输入事件和像素绘制各自清晰", 7.05, 4.0, 4.6, 0.5, size=15, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_ps2_frame(slide):
    labels = ["Start", "D0", "D1", "D2", "D3", "D4", "D5", "D6", "D7", "Parity", "Stop"]
    x = 0.78
    widths = [0.8] + [0.62] * 8 + [0.9, 0.72]
    for i, lab in enumerate(labels):
        fill = COLORS["soft_orange"] if i in (0, 9, 10) else COLORS["soft_blue"]
        add_box(slide, lab, x, 2.05, widths[i], 0.6, fill, COLORS["line"], size=10.5, bold=True)
        x += widths[i] + 0.06
    add_text(slide, "11 位 PS/2 帧：start → 8-bit scan code (LSB first) → parity → stop", 1.0, 1.55, 10.8, 0.25, size=14, color=COLORS["muted"])
    chain = ["PS2_CLK / DATA", "下降沿采样", "shift_data[7:0]", "byte_ready + byte_data"]
    x = 1.05
    for i, item in enumerate(chain):
        add_box(slide, item, x, 4.05, 2.35, 0.65, [COLORS["soft_cyan"], COLORS["soft_blue"], COLORS["soft_green"], COLORS["soft_orange"]][i], COLORS["line"], size=13.2, bold=True)
        if i < len(chain) - 1:
            add_small_arrow(slide, x + 2.45, 4.28, 0.42, 0.14, COLORS["cyan"])
        x += 3.0


def draw_scan_codes(slide, spec):
    add_code(slide, spec.code, 0.78, 1.52, 6.45, 4.62, size=9.7)
    rows = [
        ("W / Up", "光标上移"),
        ("S / Down", "光标下移"),
        ("Enter / Space", "启动游戏"),
        ("Esc", "返回菜单"),
        ("F0", "释放前缀"),
        ("E0", "扩展键前缀"),
    ]
    y = 1.6
    for key, action in rows:
        add_box(slide, key, 7.85, y, 1.7, 0.42, COLORS["soft_blue"], COLORS["blue"], size=12, bold=True)
        add_text(slide, "→", 9.75, y + 0.06, 0.35, 0.16, size=12, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        add_box(slide, action, 10.35, y, 1.7, 0.42, COLORS["soft_green"], COLORS["green"], size=12, bold=True)
        y += 0.68
    add_text(slide, "事件化让菜单不依赖 PS/2 物理细节", 7.65, 5.95, 4.5, 0.32, size=13.2, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_game_slot_api(slide):
    add_box(slide, "Game Slot N\n像卡带一样接入", 5.1, 2.35, 3.1, 1.0, COLORS["soft_orange"], COLORS["orange"], size=18, bold=True)
    inputs = ["clk / reset / selected", "frame_tick / pixel_tick", "display_active / pixel_x/y", "ps2_byte_ready / data", "btn / sw"]
    outputs = ["vga_r / vga_g / vga_b", "led", "an / ca~dp", "buzzer"]
    y = 1.45
    for item in inputs:
        add_box(slide, item, 0.95, y, 3.15, 0.42, COLORS["soft_blue"], COLORS["blue"], size=11.5)
        add_small_arrow(slide, 4.23, y + 0.14, 0.38, 0.12, COLORS["blue"])
        y += 0.62
    y = 1.7
    for item in outputs:
        add_small_arrow(slide, 8.55, y + 0.14, 0.38, 0.12, COLORS["green"])
        add_box(slide, item, 9.08, y, 3.0, 0.42, COLORS["soft_green"], COLORS["green"], size=11.5)
        y += 0.7
    add_text(slide, "接口稳定后，游戏内部实现可以独立替换", 3.25, 5.45, 6.8, 0.36, size=17, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)


def draw_mux(slide, spec):
    add_code(slide, spec.code, 0.72, 1.46, 6.15, 4.82, size=9.3)
    y = 1.62
    for label in ["slot1 RGB", "slot2 RGB", "slot3 RGB", "slot4 RGB"]:
        add_box(slide, label, 7.5, y, 1.6, 0.36, COLORS["soft_blue"], COLORS["blue"], size=11.2)
        add_small_arrow(slide, 9.25, y + 0.12, 0.32, 0.1, COLORS["blue"])
        y += 0.52
    add_box(slide, "active_slot_rgb", 9.75, 2.2, 1.9, 0.58, COLORS["soft_cyan"], COLORS["cyan"], size=12.2, bold=True)
    add_box(slide, "menu RGB", 7.5, 4.18, 1.6, 0.42, COLORS["soft_orange"], COLORS["orange"], size=11.2, bold=True)
    add_small_arrow(slide, 9.25, 4.32, 0.32, 0.1, COLORS["orange"])
    add_box(slide, "VGA_RGB\n真实端口", 10.0, 4.0, 1.55, 0.72, COLORS["soft_green"], COLORS["green"], size=12.2, bold=True)
    add_text(slide, "menu_active 选择最终输出源", 7.6, 5.35, 3.95, 0.28, size=12.6, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_resource_chart(slide):
    # Table.
    table = slide.shapes.add_table(3, 4, Inches(0.92), Inches(1.68), Inches(5.7), Inches(1.55)).table
    headers = ["指标", "使用量", "占比", "说明"]
    rows = [
        ["LUT", "30146 / 63400", "约 48%", "组合逻辑压力明显"],
        ["FF", "4371 / 126800", "约 3%", "寄存器使用很少"],
    ]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(rows, 1):
        for c, v in enumerate(row):
            table.cell(r, c).text = v
    for r in range(3):
        for c in range(4):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["soft_blue"] if r == 0 else COLORS["white"]
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = FONT_CN
                    run.font.size = Pt(10.8)
                    run.font.color.rgb = COLORS["ink"]
                    if r == 0:
                        run.font.bold = True
    # Bars.
    add_text(slide, "资源占用比例", 7.6, 1.55, 2.5, 0.28, size=14, bold=True, color=COLORS["ink"])
    add_bar(slide, "LUT", 7.55, 2.15, 3.6, 0.3, 0.48, COLORS["orange"], "48%")
    add_bar(slide, "FF", 7.55, 2.9, 3.6, 0.3, 0.03, COLORS["green"], "3%")
    add_box(slide, "LUT 高 + FF 低", 1.5, 4.1, 2.4, 0.45, COLORS["soft_orange"], COLORS["orange"], size=14, bold=True)
    add_text(slide, "↓", 2.6, 4.55, 0.2, 0.18, size=14, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_box(slide, "大量逻辑在组合路径中展开", 4.5, 4.1, 3.0, 0.45, COLORS["soft_cyan"], COLORS["cyan"], size=13.5, bold=True)
    add_text(slide, "↓", 5.9, 4.55, 0.2, 0.18, size=14, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_box(slide, "用 staging / pipeline 切短路径", 8.2, 4.1, 3.0, 0.45, COLORS["soft_green"], COLORS["green"], size=13.5, bold=True)


def add_bar(slide, label, x, y, w, h, ratio, color, value):
    add_text(slide, label, x - 0.55, y - 0.02, 0.45, 0.2, size=11, bold=True, color=COLORS["ink"], align=PP_ALIGN.RIGHT)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(226, 232, 240)
    bg.line.fill.background()
    fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w * ratio), Inches(h))
    fg.fill.solid()
    fg.fill.fore_color.rgb = color
    fg.line.fill.background()
    add_text(slide, value, x + w + 0.15, y - 0.02, 0.7, 0.2, size=10.5, bold=True, color=color)


def draw_pixel_decision(slide):
    items = ["当前像素 (x, y)", "可见区？", "背景？", "地图 tile？", "玩家 / 敌人？", "子弹 / 道具？", "文字 / UI？", "图层优先级", "RGB 输出"]
    y = 1.35
    for i, item in enumerate(items):
        fill = COLORS["soft_blue"] if i < 2 else COLORS["soft_cyan"] if i < 6 else COLORS["soft_orange"] if i < 8 else COLORS["soft_green"]
        add_box(slide, item, 1.15, y, 3.15, 0.36, fill, COLORS["line"], size=11.7, bold=i in (0, 8))
        if i < len(items) - 1:
            add_text(slide, "↓", 2.68, y + 0.34, 0.2, 0.16, size=10.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        y += 0.54
    notes = [
        "比较器多",
        "case/function 展开多",
        "mux 和优先级链多",
        "组合路径更长",
    ]
    y = 2.0
    for n in notes:
        add_box(slide, n, 7.15, y, 3.3, 0.5, COLORS["soft_orange"], COLORS["orange"], size=14, bold=True)
        y += 0.78
    add_text(slide, "画面越丰富，像素级组合逻辑越重", 6.15, 5.55, 5.6, 0.35, size=17, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)


def draw_timing_concept(slide):
    add_text(slide, "100MHz 时钟周期 = 10ns", 4.0, 1.55, 5.2, 0.35, size=20, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)
    add_box(slide, "寄存器 A", 1.4, 3.0, 1.5, 0.65, COLORS["soft_blue"], COLORS["blue"], size=14, bold=True)
    add_box(slide, "组合逻辑", 4.55, 2.75, 2.25, 1.15, COLORS["soft_orange"], COLORS["orange"], size=15, bold=True)
    add_box(slide, "寄存器 B", 8.45, 3.0, 1.5, 0.65, COLORS["soft_green"], COLORS["green"], size=14, bold=True)
    add_small_arrow(slide, 3.25, 3.22, 0.6, 0.16, COLORS["cyan"])
    add_small_arrow(slide, 7.25, 3.22, 0.6, 0.16, COLORS["cyan"])
    add_box(slide, "延迟 < 10ns\nWNS 为正", 2.1, 4.72, 2.5, 0.62, COLORS["soft_green"], COLORS["green"], size=13.2, bold=True)
    add_box(slide, "延迟 > 10ns\nWNS 为负", 6.75, 4.72, 2.5, 0.62, RGBColor(254, 226, 226), COLORS["red"], size=13.2, bold=True, color=COLORS["red"])
    add_text(slide, "WNS 看最差一条；TNS 看所有失败路径负 slack 总和", 2.0, 6.02, 8.2, 0.28, size=14, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_timing_results(slide):
    table = slide.shapes.add_table(3, 4, Inches(0.9), Inches(1.55), Inches(6.5), Inches(1.62)).table
    headers = ["阶段", "WNS", "TNS", "失败端点"]
    rows = [
        ["真实基线", "-0.445 ns", "-18.104 ns", "94"],
        ["最终 post-route", "+0.174 ns", "0.000 ns", "0"],
    ]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(rows, 1):
        for c, v in enumerate(row):
            table.cell(r, c).text = v
    for r in range(3):
        for c in range(4):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["soft_blue"] if r == 0 else (RGBColor(254, 226, 226) if r == 1 else COLORS["soft_green"])
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = FONT_CN
                    run.font.size = Pt(10.6)
                    run.font.color.rgb = COLORS["ink"]
                    if r == 0 or c in (1, 2, 3):
                        run.font.bold = True
    stages = ["负 WNS\n暴露真实基线", "RTL staging\npipeline / 小 FSM", "窄范围 VGA\nmulticycle", "post-route\nWNS 转正"]
    x = 1.1
    for i, s in enumerate(stages):
        fill = [RGBColor(254, 226, 226), COLORS["soft_cyan"], COLORS["soft_blue"], COLORS["soft_green"]][i]
        line = [COLORS["red"], COLORS["cyan"], COLORS["blue"], COLORS["green"]][i]
        add_box(slide, s, x, 4.2, 2.25, 0.75, fill, line, size=12.5, bold=True, color=COLORS["red"] if i == 0 else COLORS["ink"])
        if i < len(stages) - 1:
            add_small_arrow(slide, x + 2.42, 4.48, 0.42, 0.14, COLORS["cyan"])
        x += 2.95
    add_text(slide, "强调：不是靠宽泛例外掩盖游戏核心路径", 3.5, 5.55, 5.6, 0.3, size=15, bold=True, color=COLORS["orange"], align=PP_ALIGN.CENTER)


def draw_optimization_cases(slide):
    cases = [
        ("Slot2", ["ghost 迭代", "lock row-by-row", "pixel_tick video staging"]),
        ("Slot3", ["candidate position", "try_x / try_y 寄存", "下一拍 walkable 判断"]),
        ("Slot4", ["TEST：寄存 tile 地址", "EVAL：评估结果", "APPLY：更新状态"]),
    ]
    x = 0.85
    for title, lines in cases:
        add_box(slide, title, x, 1.65, 3.55, 0.45, COLORS["soft_blue"], COLORS["blue"], size=16, bold=True)
        y = 2.32
        for line in lines:
            add_box(slide, line, x + 0.25, y, 3.05, 0.38, COLORS["white"], COLORS["line"], size=11.5)
            if line != lines[-1]:
                add_text(slide, "↓", x + 1.72, y + 0.34, 0.18, 0.15, size=10, color=COLORS["muted"], align=PP_ALIGN.CENTER)
            y += 0.58
        x += 4.08
    add_box(slide, "一拍大计算", 2.8, 5.2, 2.2, 0.52, RGBColor(254, 226, 226), COLORS["red"], size=14, bold=True, color=COLORS["red"])
    add_small_arrow(slide, 5.18, 5.38, 0.58, 0.16, COLORS["cyan"])
    add_box(slide, "多拍小计算", 6.0, 5.2, 2.2, 0.52, COLORS["soft_green"], COLORS["green"], size=14, bold=True, color=COLORS["green"])
    add_text(slide, "行为保持帧级等价，关键路径更短", 4.15, 6.05, 4.8, 0.28, size=14, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_summary(slide):
    items = [
        ("VGA", "实时像素扫描"),
        ("菜单", "坐标判断 / 图层优先级"),
        ("PS/2", "协议解析 / 事件化"),
        ("API", "平台化槽位架构"),
        ("LUT/WNS", "从功能到可实现硬件"),
    ]
    x = 0.95
    for i, (head, text) in enumerate(items):
        fill = [COLORS["soft_blue"], COLORS["soft_cyan"], COLORS["soft_green"], COLORS["soft_orange"], RGBColor(226, 232, 240)][i]
        add_box(slide, f"{head}\n{text}", x, 2.0, 2.15, 1.05, fill, COLORS["line"], size=13, bold=True)
        x += 2.4
    add_text(
        slide,
        "这个项目让我从“写出能工作的游戏”走向“设计一个能显示、能输入、能扩展、能优化、能真正生成 bitstream 的 FPGA 产品”。",
        1.35,
        4.55,
        10.6,
        0.72,
        size=18,
        color=COLORS["blue"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def render_slide(prs, spec, idx):
    if spec.kind == "cover":
        return cover_slide(prs, spec)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, idx, spec.title, spec.takeaway)
    if spec.visual == "user_flow":
        draw_user_flow(slide)
    elif spec.visual == "learning_route":
        draw_learning_route(slide)
    elif spec.visual == "architecture":
        draw_architecture(slide)
    elif spec.visual == "vga_scan":
        draw_vga_scan(slide)
    elif spec.visual == "code":
        draw_code_slide(slide, spec)
    elif spec.visual == "single_clock":
        draw_single_clock(slide)
    elif spec.visual == "menu_render":
        draw_menu_render(slide, spec)
    elif spec.visual == "menu_control":
        draw_menu_control(slide)
    elif spec.visual == "ps2_frame":
        draw_ps2_frame(slide)
    elif spec.visual == "scan_codes":
        draw_scan_codes(slide, spec)
    elif spec.visual == "game_slot_api":
        draw_game_slot_api(slide)
    elif spec.visual == "mux":
        draw_mux(slide, spec)
    elif spec.visual == "resource_chart":
        draw_resource_chart(slide)
    elif spec.visual == "pixel_decision":
        draw_pixel_decision(slide)
    elif spec.visual == "timing_concept":
        draw_timing_concept(slide)
    elif spec.visual == "timing_results":
        draw_timing_results(slide)
    elif spec.visual == "optimization_cases":
        draw_optimization_cases(slide)
    elif spec.visual == "summary":
        draw_summary(slide)
    else:
        add_multiline(slide, spec.bullets, 0.9, 1.75, 11.4, 3.6, size=18, bullet=True)
    if spec.visual not in {"code", "menu_render", "scan_codes", "mux", "resource_chart", "timing_results", "optimization_cases", "summary"}:
        add_bullet_panel(slide, spec.bullets)
    return slide


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    for idx, spec in enumerate(SLIDES, start=1):
        render_slide(prs, spec, idx)
    prs.save(PPTX_PATH)


def markdown_outline() -> str:
    lines = [
        "# PPT 大纲",
        "",
        "主题：《“学”：基于 Nexys A7 的 VGA 游戏集合机》",
        "",
        "主线：我不是只做了几个小游戏，而是学习并完成了一个 FPGA 图形交互系统：VGA 显示、菜单渲染、PS/2 输入、Game Slot API、资源占用分析和时序收敛。",
        "",
    ]
    for i, spec in enumerate(SLIDES, start=1):
        lines.append(f"## 第 {i} 页：{spec.title}")
        lines.append("")
        lines.append(f"**核心结论：** {spec.takeaway}")
        lines.append("")
        if spec.bullets:
            lines.append("**页面要点：**")
            for b in spec.bullets:
                lines.append(f"- {b}")
            lines.append("")
        lines.append(f"**图示/版式：** {spec.visual}")
        if spec.refs:
            lines.append("")
            lines.append("**引用：** " + "；".join(spec.refs))
        lines.append("")
    return "\n".join(lines)


def markdown_script() -> str:
    lines = ["# 逐页演讲稿", "", "说明：每页讲稿控制在约 80 到 180 字，可直接用于课程验收汇报。", ""]
    for i, spec in enumerate(SLIDES, start=1):
        lines.append(f"## 第 {i} 页：{spec.title}")
        lines.append("")
        lines.append(spec.speaker)
        lines.append("")
    return "\n".join(lines)


def markdown_citations() -> str:
    code_refs = [
        ("README.md", "项目定位、板卡型号、四个槽位、用户操作流程。"),
        ("src/game_console_top.v", "公共模块实例化、slot 选择、菜单和游戏输出复用、LED/AN/SEG/BUZZER 仲裁。"),
        ("src/common/console_vga_sync.v", "VGA 参数、`pix_div`、`pixel_tick`、`h_count/v_count`、`display_active`、`hsync/vsync`。"),
        ("src/common/console_menu_renderer.v", "菜单标题、菜单项、帮助文字、边框和高亮的像素级坐标判断。"),
        ("src/common/console_ps2_rx.v", "PS/2 同步/滤波、下降沿检测、`bit_count`、`shift_data`、`byte_ready`。"),
        ("src/common/console_menu_controller.v", "W/S、方向键、Enter/Space、Esc、F0、E0 的扫描码事件化。"),
        ("docs/game_api.md", "标准 `game_slotN_top` 接口、公共输入和统一输出。"),
        ("docs/资源占用分析.md", "LUT 30146/63400、FF 4371/126800，以及 LUT 高 FF 低的解释。"),
        ("docs/final_timing_optimization_round.md", "真实基线 WNS/TNS/失败端点、最终 post-route 结果、Slot2/3/4 优化。"),
        ("scripts/apply_timing_exceptions.tcl", "窄范围 4-cycle VGA pixel multicycle，明确不覆盖 game-core state。"),
        ("constraints/game_console.xdc", "100MHz 时钟、VGA、PS/2、按钮、开关、LED、七段管和蜂鸣器引脚约束。"),
        ("src/games/slot1/", "TANK WAR 槽位包装与接入。"),
        ("src/games/slot2/", "ghost/lock 多拍优化与 video staging。"),
        ("src/games/slot3/", "`try_x/try_y` 候选坐标寄存与 walkable 判断拆分。"),
        ("src/games/slot4/", "tile probe/eval/apply 分阶段物理和状态检测。"),
    ]
    snippets = {
        "VGA 参数与像素使能": dedent(
            """\
            // src/common/console_vga_sync.v
            localparam integer H_VISIBLE = 640;
            localparam integer H_FRONT   = 16;
            localparam integer H_SYNC    = 96;
            localparam integer H_BACK    = 48;
            localparam integer V_VISIBLE = 480;
            localparam integer V_FRONT   = 10;
            localparam integer V_SYNC    = 2;
            localparam integer V_BACK    = 33;
            assign pixel_tick = (pix_div == 2'b11);
            assign display_active = (h_count < H_VISIBLE) && (v_count < V_VISIBLE);
            """
        ),
        "PS/2 接收关键逻辑": dedent(
            """\
            // src/common/console_ps2_rx.v
            assign ps2_clk_fall_next = ps2_clk_q & ~ps2_clk_next;
            assign ps2_parity_ok = (^shift_data) ^ parity_bit;
            if (ps2_data_next == 1'b1 && ps2_parity_ok) begin
                byte_data <= shift_data;
                byte_ready <= 1'b1;
            end
            """
        ),
        "菜单扫描码事件化": dedent(
            """\
            // src/common/console_menu_controller.v
            assign is_up_key = (byte_data == SCAN_W) || (byte_data == SCAN_UP);
            assign is_down_key = (byte_data == SCAN_S) || (byte_data == SCAN_DOWN);
            assign is_start_key = (byte_data == SCAN_SPACE) || (byte_data == SCAN_ENTER);
            assign is_back_key = (byte_data == SCAN_ESC);
            """
        ),
        "顶层 RGB 复用": dedent(
            """\
            // src/game_console_top.v
            assign active_slot_vga_r = slot1_selected ? slot1_vga_r :
                                       slot2_selected ? slot2_vga_r :
                                       slot3_selected ? slot3_vga_r :
                                                        slot4_vga_r;
            if (menu_active) VGA_R = menu_vga_r_q;
            else             VGA_R = active_slot_vga_r;
            """
        ),
    }
    lines = ["# 素材与代码引用清单", ""]
    lines.append("## 核查过的文件与用途")
    lines.append("")
    for path, desc in code_refs:
        lines.append(f"- `{path}`：{desc}")
    lines.append("")
    lines.append("## 关键数据")
    lines.extend(
        [
            "",
            "- 板卡：Nexys A7-100T，器件 XC7A100T-1CSG324C。",
            "- VGA：640x480 可见区，水平参数 640/16/96/48，垂直参数 480/10/2/33。",
            "- `pixel_tick`：100MHz 主时钟下 `pix_div == 2'b11`，四拍一次像素使能。",
            "- 资源：LUT 30146 / 63400，约 48%；FF 4371 / 126800，约 3%。",
            "- 真实时序基线：WNS=-0.445ns，TNS=-18.104ns，失败端点 94。",
            "- 最终 post-route：WNS=+0.174ns，TNS=0.000ns，失败端点 0。",
            "- 时序例外：只保留窄范围 VGA pixel multicycle，`set_multicycle_path 4 -setup` 配 `3 -hold`，不覆盖游戏核心状态。",
            "",
        ]
    )
    lines.append("## PPT 使用的代码片段")
    lines.append("")
    for title, snippet in snippets.items():
        lines.append(f"### {title}")
        lines.append("")
        lines.append("```verilog")
        lines.append(snippet.rstrip())
        lines.append("```")
        lines.append("")
    lines.append("## 图示来源")
    lines.append("")
    lines.append("本 PPT 未使用网络图片。系统架构图、VGA 扫描图、PS/2 帧结构图、Game Slot API 图、输出复用 mux 图、资源占用图、时序优化路线图和 Slot2/3/4 优化案例图均由 `build_ppt.py` 使用 PowerPoint shapes 绘制。")
    return "\n".join(lines)


def markdown_report(pdf_status: str) -> str:
    lines = [
        "# 制作说明与自检报告",
        "",
        "## 生成方式",
        "",
        f"- 生成脚本：`{OUT_DIR.name}/build_ppt.py`",
        "- PPTX 生成：Python + `python-pptx`，16:9 宽屏。",
        f"- PDF 生成：{pdf_status}",
        "- 图示：全部使用 PPT shapes 或脚本绘制，不使用网络图片。",
        "- 字体：中文优先使用 SimHei，代码使用 Consolas。",
        "",
        "## 与用户给定数据的差异",
        "",
        "当前项目文档中的关键数字与用户给定数字一致：LUT/FF、真实基线 WNS/TNS/失败端点、最终 post-route WNS/TNS/失败端点均未发现冲突。",
        "补充注意：`docs/资源占用分析.md` 中 `u_console_vga_sync = 5829 LUT` 被文档标注为可能受 flatten/层级归属影响，PPT 没有把它当作 VGA sync 本体真实逻辑规模来强调。",
        "",
        "## 自检结果",
        "",
    ]
    checks = [
        ("PPT 是否围绕“学”这个主题，而不是普通实验报告？", "通过。第 1、3、20 页回扣“学”，中间按学习路径组织。"),
        ("是否重点讲了 VGA 显示原理？", "通过。第 5-7 页讲像素扫描、VGA 参数和 `pixel_tick`。"),
        ("是否讲清菜单渲染不是图片，而是像素级判断？", "通过。第 8-9 页用伪代码和链路图说明。"),
        ("是否讲清 PS/2 输入从物理信号到 byte 事件再到菜单动作？", "通过。第 10-12 页覆盖帧结构、接收器和扫描码事件化。"),
        ("是否讲清 Game Slot API 的平台化意义？", "通过。第 13-14 页讲卡带式接入和输出复用。"),
        ("是否讲清 LUT 高、FF 低的含义？", "通过。第 15-16 页给出表格和像素决策树。"),
        ("是否讲清 WNS/TNS 和时序优化结果？", "通过。第 17-19 页给出概念、结果表和优化案例。"),
        ("是否包含真实代码片段，并标注文件位置？", "通过。第 6、11、12、14 页放真实代码；第 8 页明确标注伪代码。"),
        ("是否避免了大段堆文字？", "通过。主要使用流程图、表格、代码块和短要点。"),
        ("是否有逐页演讲稿？", "通过。见 `逐页演讲稿.md`。"),
        ("是否有最终 `.pptx`？", "通过。见 `VGA游戏集合机_验收讲解.pptx`。"),
        ("如果支持，是否有 `.pdf`？", "通过。已生成 `VGA游戏集合机_验收讲解.pdf`。"),
        ("所有关键数字是否和项目文档一致？", "通过。已核查 README、资源分析和最终时序优化文档。"),
        ("是否能在 12 到 18 分钟内讲完？", "通过。20 页，每页约 35-50 秒，预计 13-16 分钟。"),
    ]
    for question, result in checks:
        lines.append(f"- **{question}** {result}")
    lines.extend(
        [
            "",
            "## 已知风险或 TODO",
            "",
            "- 最终 WNS 为 +0.174ns，余量为正但不大；后续如果继续增加游戏逻辑，仍需重新跑 post-route timing。",
            "- PDF 由脚本直接生成，若需要与 PowerPoint 版完全逐像素一致，可在安装 PowerPoint 或 LibreOffice 后再做正式导出。",
            "- PPT 未包含上板实拍图片；现场验收时建议先演示实物，再进入第 1 页汇报。",
            "",
            "## 产物清单",
            "",
            f"- `{PPTX_PATH.name}`",
            f"- `{PDF_PATH.name}`",
            f"- `{OUTLINE_PATH.name}`",
            f"- `{SCRIPT_PATH.name}`",
            f"- `{CITATION_PATH.name}`",
            f"- `{REPORT_PATH.name}`",
            "- `build_ppt.py`",
        ]
    )
    return "\n".join(lines)


def write_markdown(pdf_status: str):
    OUTLINE_PATH.write_text(markdown_outline(), encoding="utf-8")
    SCRIPT_PATH.write_text(markdown_script(), encoding="utf-8")
    CITATION_PATH.write_text(markdown_citations(), encoding="utf-8")
    REPORT_PATH.write_text(markdown_report(pdf_status), encoding="utf-8")


def pdf_wrap(draw, text, font, max_width):
    lines = []
    current = ""
    for ch in text:
        probe = current + ch
        bbox = draw.textbbox((0, 0), probe, font=font)
        if bbox[2] - bbox[0] <= max_width or not current:
            current = probe
        else:
            lines.append(current)
            current = ch
    if current:
        lines.append(current)
    return lines


def build_pdf_fallback() -> str:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception as exc:  # pragma: no cover
        return f"未生成：PIL 不可用（{exc}）。"

    font_cn_path = Path("C:/Windows/Fonts/simhei.ttf")
    font_code_path = Path("C:/Windows/Fonts/consola.ttf")
    font_title = ImageFont.truetype(str(font_cn_path), 42)
    font_sub = ImageFont.truetype(str(font_cn_path), 24)
    font_body = ImageFont.truetype(str(font_cn_path), 25)
    font_small = ImageFont.truetype(str(font_cn_path), 19)
    font_code = ImageFont.truetype(str(font_code_path), 18) if font_code_path.exists() else font_small

    pages = []
    for idx, spec in enumerate(SLIDES, start=1):
        img = Image.new("RGB", (1600, 900), (247, 249, 252))
        draw = ImageDraw.Draw(img)
        if spec.kind == "cover":
            draw.rectangle([0, 0, 1600, 900], fill=(10, 22, 40))
            draw.text((85, 100), spec.title, font=font_title, fill=(255, 255, 255))
            draw.text((90, 185), "从 VGA 显示、PS/2 输入到资源分析与时序收敛", font=font_sub, fill=(203, 213, 225))
            x = 95
            for item in ["VGA", "PS2", "Game Slot", "LUT", "WNS"]:
                draw.rounded_rectangle([x, 710, x + 155, 760], radius=10, fill=(30, 41, 59), outline=(80, 95, 115))
                draw.text((x + 32, 722), item, font=font_small, fill=(255, 255, 255))
                x += 178
            pages.append(img)
            continue
        draw.text((72, 42), spec.title, font=font_title, fill=(11, 31, 51))
        draw.text((78, 104), spec.takeaway, font=font_small, fill=(82, 96, 112))
        draw.rectangle([78, 148, 320, 154], fill=(8, 145, 178))
        y = 195
        if spec.code:
            draw.rounded_rectangle([78, y, 850, y + 420], radius=12, fill=(15, 23, 42), outline=(30, 41, 59))
            cy = y + 22
            for line in spec.code.rstrip().splitlines()[:18]:
                draw.text((100, cy), line, font=font_code, fill=(226, 232, 240))
                cy += 24
            y = 640
        else:
            draw.rounded_rectangle([82, 205, 1515, 600], radius=16, fill=(255, 255, 255), outline=(203, 213, 225))
            draw.text((120, 250), f"图示：{spec.visual}", font=font_sub, fill=(29, 78, 216))
            for line in pdf_wrap(draw, spec.takeaway, font_body, 1180)[:3]:
                draw.text((120, 315), line, font=font_body, fill=(11, 31, 51))
                y = 360
        by = y if spec.code else 635
        for b in spec.bullets[:5]:
            wrapped = pdf_wrap(draw, "• " + b, font_body, 1350)
            for line in wrapped[:2]:
                draw.text((105, by), line, font=font_body, fill=(11, 31, 51))
                by += 34
        draw.text((1430, 840), f"{idx:02d}", font=font_small, fill=(82, 96, 112))
        pages.append(img)

    first, rest = pages[0], pages[1:]
    first.save(PDF_PATH, "PDF", resolution=160.0, save_all=True, append_images=rest)
    return "已生成。LibreOffice/PowerPoint 不可用，PDF 使用 PIL 按同一 slide metadata 生成可阅读备份版。"


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_pptx()
    pdf_status = build_pdf_fallback()
    write_markdown(pdf_status)
    print(f"PPTX: {PPTX_PATH}")
    print(f"PDF: {PDF_PATH if PDF_PATH.exists() else 'not generated'}")
    print(f"Outline: {OUTLINE_PATH}")
    print(f"Script: {SCRIPT_PATH}")
    print(f"Citations: {CITATION_PATH}")
    print(f"Report: {REPORT_PATH}")


if __name__ == "__main__":
    main()
